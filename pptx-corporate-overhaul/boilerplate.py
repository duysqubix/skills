"""
PowerPoint Corporate Overhaul — Boilerplate
============================================
Copy this file as your starting scaffold when overhauling a PowerPoint.
Adapt colors, data, and slide structure to match the source deck.

Dependencies: python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree

# ═══════════════════════════════════════════════════════════════════
# CHOOSE ONE PALETTE — comment out the other
# ═══════════════════════════════════════════════════════════════════

# ─── Corporate Dark ───────────────────────────────────────────────
NAVY        = RGBColor(0x0F, 0x1A, 0x2E)
NAVY_2      = RGBColor(0x17, 0x25, 0x3D)
CARD        = RGBColor(0x1B, 0x2A, 0x44)
CARD_HI     = RGBColor(0x22, 0x33, 0x52)
BORDER      = RGBColor(0x2E, 0x40, 0x60)
RED         = RGBColor(0xC8, 0x10, 0x2E)
RED_DARK    = RGBColor(0x8B, 0x0B, 0x20)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
BLUE_SOFT   = RGBColor(0x60, 0xA5, 0xFA)
WHITE       = RGBColor(0xF3, 0xF4, 0xF6)
WHITE_DIM   = RGBColor(0xCB, 0xD5, 0xE1)
SLATE       = RGBColor(0x94, 0xA3, 0xB8)
SLATE_DK    = RGBColor(0x64, 0x74, 0x8B)
LOGO_PATH   = "logo_white.png"   # inverted logo for dark bg

# ─── Corporate Light (uncomment to use) ───────────────────────────
# NAVY        = RGBColor(0xF7, 0xF9, 0xFC)
# NAVY_2      = RGBColor(0xEE, 0xF2, 0xF7)
# CARD        = RGBColor(0xFF, 0xFF, 0xFF)
# CARD_HI     = RGBColor(0xF1, 0xF5, 0xF9)
# BORDER      = RGBColor(0xD9, 0xE0, 0xE9)
# RED         = RGBColor(0xC8, 0x10, 0x2E)
# RED_DARK    = RGBColor(0x8B, 0x0B, 0x20)
# AMBER       = RGBColor(0xB4, 0x5B, 0x0A)
# GREEN       = RGBColor(0x04, 0x7A, 0x56)
# BLUE_SOFT   = RGBColor(0x1D, 0x4E, 0xD8)
# WHITE       = RGBColor(0x0B, 0x13, 0x26)
# WHITE_DIM   = RGBColor(0x2D, 0x3B, 0x55)
# SLATE       = RGBColor(0x5B, 0x6B, 0x82)
# SLATE_DK    = RGBColor(0x9B, 0xA8, 0xBC)
# LOGO_PATH   = "logo.png"         # original dark logo for light bg


# ═══════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def set_line(shape, color, weight=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(weight)

def add_rect(slide, x, y, w, h, fill=CARD, line=None, corner=False):
    """Add a rectangle or rounded rectangle. corner=True applies slight rounding."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if corner:
        sh.adjustments[0] = 0.06
    set_fill(sh, fill)
    if line is None:
        no_line(sh)
    else:
        set_line(sh, line, 0.75)
    sh.shadow.inherit = False
    return sh

def add_text(slide, x, y, w, h, text, *, font_size=14, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="Calibri", italic=False, spacing=None):
    """Add a single-style text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font_name
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def add_multi_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP, spacing=None):
    """Add a text box with mixed-style runs.
    runs = list of dicts: {text, size, bold, color, font, italic, newline}
    newline=True starts a new paragraph.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    cur = tf.paragraphs[0]
    cur.alignment = align
    if spacing is not None:
        cur.line_spacing = spacing
    first = True
    for run in runs:
        if run.get("newline") and not first:
            cur = tf.add_paragraph()
            cur.alignment = run.get("align", align)
            if spacing is not None:
                cur.line_spacing = spacing
        r = cur.add_run()
        r.text = run["text"]
        r.font.name = run.get("font", "Calibri")
        r.font.size = Pt(run.get("size", 14))
        r.font.bold = run.get("bold", False)
        r.font.italic = run.get("italic", False)
        r.font.color.rgb = run.get("color", WHITE)
        first = False
    return tb

def set_slide_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_chrome(slide, title, subtitle=None, slide_num=None, total=9,
                     footer_text="COMPANY NAME  ·  REPORT TYPE  ·  PERIOD"):
    """Standard header bar, logo, footer line, footer text, and slide counter."""
    add_rect(slide, Inches(0), Inches(0.48), Inches(13.333), Inches(0.04), fill=RED)
    add_text(slide, Inches(0.55), Inches(0.10), Inches(9.5), Inches(0.45),
             title, font_size=24, bold=True, color=WHITE,
             font_name="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.04), Inches(9.5), Inches(0.18),
                 subtitle.upper(), font_size=9, bold=True, color=RED,
                 font_name="Calibri", align=PP_ALIGN.LEFT)
    slide.shapes.add_picture(LOGO_PATH, Inches(11.8), Inches(0.12), height=Inches(0.34))
    add_rect(slide, Inches(0.55), Inches(7.10), Inches(12.23), Inches(0.015), fill=BORDER)
    add_text(slide, Inches(0.55), Inches(7.18), Inches(8), Inches(0.24),
             footer_text, font_size=9, color=SLATE, bold=True, font_name="Calibri")

    if slide_num is not None:
        add_text(slide, Inches(11.5), Inches(7.18), Inches(1.28), Inches(0.24),
                 f"{slide_num:02d} / {total:02d}",
                 font_size=9, color=SLATE, bold=True, font_name="Calibri",
                 align=PP_ALIGN.RIGHT)


def add_kpi_card(slide, x, y, w, h, value, label, *, value_color=WHITE,
                 accent=RED, bg=CARD, sublabel=None):
    """KPI chip with left accent stripe, big value, and label."""
    add_rect(slide, x, y, w, h, fill=bg, corner=True)
    add_rect(slide, x, y, Inches(0.08), h, fill=accent)
    add_text(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.4), Inches(0.75),
             value, font_size=40, bold=True, color=value_color,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.25), y + h - Inches(0.45), w - Inches(0.4), Inches(0.26),
             label.upper(), font_size=9, bold=True, color=SLATE,
             align=PP_ALIGN.LEFT)
    if sublabel:
        add_text(slide, x + Inches(0.25), y + h - Inches(0.22), w - Inches(0.4), Inches(0.2),
                 sublabel, font_size=9, color=WHITE_DIM,
                 align=PP_ALIGN.LEFT, italic=True)


def add_leaderboard_row(slide, x, y, w, h, rank, name, subtitle, count, clr,
                        max_val, bar_max_w):
    """Single leaderboard row with progress bar."""
    add_rect(slide, x, y, w, h, fill=CARD, corner=True)
    add_text(slide, x + Inches(0.2), y, Inches(0.6), h,
             rank, font_size=22, bold=True, color=SLATE_DK,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.85), y + Inches(0.12), Inches(3.8), Inches(0.32),
             name, font_size=14, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.85), y + Inches(0.43), Inches(4.3), Inches(0.32),
             subtitle, font_size=8.5, bold=True, color=SLATE)
    add_text(slide, x + Inches(4.4), y, Inches(0.8), h,
             str(count), font_size=26, bold=True, color=clr,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    bar_x = x + Inches(5.25)
    bar_y = y + Inches(0.28)
    bar_h = Inches(0.32)
    add_rect(slide, bar_x, bar_y, bar_max_w, bar_h, fill=CARD_HI, corner=True)
    fill_w = Emu(int(bar_max_w * (count / max_val)))
    if fill_w > 0:
        add_rect(slide, bar_x, bar_y, fill_w, bar_h, fill=clr, corner=True)


def style_chart(chart, series_color, data_label_fmt='#,##0', show_legend=False):
    """Apply corporate styling to a column chart."""
    chart.has_title = False
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.color.rgb = WHITE_DIM
        chart.legend.font.size = Pt(11)
    chart.category_axis.tick_labels.font.color.rgb = WHITE_DIM
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.color.rgb = WHITE_DIM
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.number_format = data_label_fmt
    chart.category_axis.format.line.color.rgb = BORDER
    chart.value_axis.format.line.fill.background()
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
    for ser in chart.series:
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = series_color
        ser.data_labels.show_value = True
        ser.data_labels.font.color.rgb = WHITE
        ser.data_labels.font.size = Pt(10)
        ser.data_labels.font.bold = True
        ser.data_labels.number_format = data_label_fmt
        ser.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


# ═══════════════════════════════════════════════════════════════════
# MINIMAL WORKING EXAMPLE
# ═══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(6.5), Inches(7.5), fill=NAVY)
    add_text(s, Inches(0.7), Inches(2.1), Inches(5.8), Inches(2.5),
             "Q1 REVIEW\n2026", font_size=60, bold=True, color=WHITE)
    add_rect(s, Inches(0.7), Inches(6.05), Inches(0.8), Inches(0.03), fill=RED)

    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    add_slide_chrome(s, "Safety Performance",
                     subtitle="INCIDENT TRENDING · Q1 2026", slide_num=2)

    add_kpi_card(s, Inches(0.55), Inches(0.95), Inches(4.0), Inches(1.10),
                 "0", "LOST TIME INCIDENTS", accent=GREEN,
                 sublabel="Year-to-date · Q1 2026")
    add_kpi_card(s, Inches(4.65), Inches(0.95), Inches(4.0), Inches(1.10),
                 "0", "OSHA RECORDABLES", accent=GREEN,
                 sublabel="Year-to-date · Q1 2026")
    add_kpi_card(s, Inches(8.75), Inches(0.95), Inches(4.0), Inches(1.10),
                 "0", "VEHICLE INCIDENTS", accent=GREEN,
                 sublabel="Year-to-date · Q1 2026")

    add_rect(s, Inches(0.55), Inches(2.25), Inches(12.23), Inches(4.7),
             fill=CARD, corner=True)

    cd = CategoryChartData()
    cd.categories = ["Jan", "Feb", "Mar"]
    cd.add_series("Incidents", (0, 0, 0))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(2.6), Inches(11.3), Inches(4.0), cd
    ).chart
    style_chart(chart, GREEN, data_label_fmt='0')

    prs.save("corporate_overhaul_example.pptx")
    print("Saved: corporate_overhaul_example.pptx")
